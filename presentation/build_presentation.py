#!/usr/bin/env python3
"""Build CIC-IDS ablation deck. Run: PYTHONPATH=../.pptx_deps python build_presentation.py"""
import sys
from pathlib import Path

root = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(root / ".pptx_deps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title: str, subtitle: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2))
        tx.text_frame.text = title
        for p in tx.text_frame.paragraphs:
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(20, 60, 120)
        tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(2))
        tx2.text_frame.text = subtitle
        for p in tx2.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(60, 60, 60)

    def add_bullet_slide(title: str, bullets: list):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.8))
        t.text_frame.text = title
        for p in t.text_frame.paragraphs:
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(20, 60, 120)
        box = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.5), Inches(5.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            par.text = b
            par.level = 0
            par.font.size = Pt(18)
            par.space_after = Pt(10)

    def add_table_slide(title: str, headers: list, rows: list):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.75))
        t.text_frame.text = title
        for p in t.text_frame.paragraphs:
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = RGBColor(20, 60, 120)
        nc = len(headers)
        nr = 1 + len(rows)
        tbl = slide.shapes.add_table(nr, nc, Inches(0.7), Inches(1.35), Inches(12), Inches(0.45 * nr)).table
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(14)
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                c = tbl.cell(i + 1, j)
                c.text = str(val)
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(13)

    add_title_slide(
        "Network Intrusion Detection",
        "Hybrid ML + Deep Learning on CIC-IDS 2017\nAML-DL Combined — May 2026",
    )

    add_bullet_slide("Problem & Goal", [
        "Goal: flag malicious flows when supervised labels are limited.",
        "Setting: unsupervised anomaly detection — fit on benign (BENIGN) traffic only.",
        "Compare Model A (OCSVM), Model B (MAE reconstruction), Model C (MAE + OCSVM hybrid).",
    ])

    add_bullet_slide("Dataset: CIC-IDS 2017", [
        "Eight CSV segments merged → data/raw/cicids2017_all.csv (~2.83 million flows).",
        "Binary task: BENIGN vs any attack (multi-class labels collapsed).",
        "Features: 77 dims after preprocessing + engineering; identifiers excluded.",
        "Excluded: Flow ID, Source/Destination IP & ports, Timestamp.",
    ])

    add_bullet_slide("Models & Experimental Setup", [
        "Model A: One-Class SVM (RBF, ν=0.01) on scaled tabular features.",
        "Model B: Tabular MAE — threshold on reconstruction MSE vs benign median + 2σ.",
        "Model C: Hybrid — OCSVM on frozen MAE embeddings.",
        "Evaluation: python scripts/run_ablation.py --fast → 5k train / 5k test, seed 42.",
    ])

    add_bullet_slide("MAE / Training Notes", [
        "reports/mae_training_history.csv: logged multi-epoch run from earlier train_mae.py session.",
        "CIC-IDS ablation used a 2-epoch MAE warm-up on benign subset (77 input dims).",
        "MAE config: hidden 128, 4 layers, 4 heads, mask_ratio 0.15 (see config/config.yaml).",
    ])

    add_table_slide(
        "Results — All Metrics (from results/ablation_table.csv)",
        ["Model", "Accuracy", "Precision", "Recall", "F1", "ROC-AUC"],
        [
            ["A: OCSVM only", "0.8624", "0.7602", "0.4363", "0.5544", "0.7943"],
            ["B: MAE only", "0.8024", "0.1111", "0.0010", "0.0020", "0.7059"],
            ["C: Hybrid", "0.7950", "0.0769", "0.0041", "0.0077", "0.3763"],
        ],
    )

    add_bullet_slide("Discussion", [
        "Strongest on this run: Model A (ROC-AUC 0.794, F1 0.554).",
        "B/C: MAE barely trained for CIC in this session → very low recall and F1.",
        "Honest takeaway: deep models need full benign pretraining to compare fairly.",
        "Artifacts: results/ablation_table.csv, models/mae_pretrained.pt.",
    ])

    add_title_slide("Thank you", "Q&A\nRepo: AML-DL-COMBINED")

    out_dir = root / "presentation"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "CIC-IDS_ablation_presentation.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
